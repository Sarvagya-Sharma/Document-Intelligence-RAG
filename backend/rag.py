from http import client
import re
import io
import time
import contextlib
import chromadb
import pymupdf
import pytesseract
import hashlib
import os
import numpy as np
from PIL import Image
from tqdm import tqdm
from dotenv import load_dotenv
from docling.document_converter import DocumentConverter
from sentence_transformers import SentenceTransformer, CrossEncoder
from rank_bm25 import BM25Okapi
from google import genai
from google.genai import types
import json
from langsmith import Client, traceable
import openpyxl
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches
import requests

load_dotenv()
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
genai_client = genai.Client()
chroma_client = chromadb.Client()

EMBED_MODEL = SentenceTransformer("all-MiniLM-L6-v2")
EMBED_MODEL_NAME = "all-MiniLM-L6-v2"
RERANK_MODEL = CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2")

LLM_MODEL = "gemma-4-26b-a4b-it"
OLLAMA_MODEL = "mistral"
OLLAMA_URL   = "http://localhost:11434/api/generate"
COLLECTION_NAME = "Project_doc"
SEMANTIC_BREAKPOINT_THRESHOLD = 0.35   
MAX_CHUNK_TOKENS = 800               
MIN_CHUNK_CHARS  = 80                
DENSE_TOP_K  = 20  
SPARSE_TOP_K = 20
FINAL_TOP_K  = 5   
RRF_K        = 60   
BM25_WEIGHT  = 0.4

DEFAULT_TOP_K = FINAL_TOP_K


def text_format(text: str) -> str:
    return text.replace("\n", " ").strip()


def cosine_distance(a: np.ndarray, b: np.ndarray) -> float:
    a, b = np.array(a), np.array(b)
    norm = np.linalg.norm(a) * np.linalg.norm(b)
    if norm == 0:
        return 1.0
    return float(1 - np.dot(a, b) / norm)


def open_and_read_pdf(pdf_path: str):
    texts = []
    with pymupdf.open(pdf_path) as doc:
        for i, page in tqdm(enumerate(doc), total=len(doc)):
            text = page.get_text().strip()
            if len(text) < 50:
                pix = page.get_pixmap(dpi=300)
                img = Image.open(io.BytesIO(pix.tobytes()))
                text = pytesseract.image_to_string(img)
            texts.append({
                "source": pdf_path,
                "source_page": f"Page {i + 1}",   # ← string
                "text": text_format(text)
            })
    return texts


def open_and_read_document(file_path: str):
    converter = DocumentConverter()
    result = converter.convert(file_path)
    full_text = text_format(result.document.export_to_markdown())
    return [{"source": file_path, "source_page": "Page 1", "text": full_text}]  # ← string


def open_and_read_excel(file_path: str) -> list[dict]:
    wb = openpyxl.load_workbook(file_path, data_only=True)
    pages = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not any(any(cell is not None for cell in row) for row in rows):
            continue
        lines = []
        for i, row in enumerate(rows):
            cells = [str(c) if c is not None else "" for c in row]
            lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                lines.append("|" + "|".join(["---"] * len(cells)) + "|")
        text = f"Sheet: {sheet_name}\n\n" + "\n".join(lines)
        pages.append({
            "source": file_path,
            "source_page": f"Sheet: {sheet_name}",   # ← string
            "text": text_format(text)
        })
    return pages

def _extract_shape_text(shape) -> list[str]:
    parts = []

    if shape.shape_type == 6:  
        for child in shape.shapes:
            parts.extend(_extract_shape_text(child))
        return parts


    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            parts.append(text)

    elif shape.has_table:
        rows_text = []
        for i, row in enumerate(shape.table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            rows_text.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows_text.append("|" + "|".join(["---"] * len(cells)) + "|")
        if rows_text:
            parts.append("\n".join(rows_text))

    return parts


def open_and_read_pptx(file_path: str) -> list[dict]:
    prs = Presentation(file_path)
    pages = []
    for slide_num, slide in enumerate(tqdm(prs.slides, desc="Reading slides")):
        parts = []
        for shape in slide.shapes:
            parts.extend(_extract_shape_text(shape))
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f"[Speaker Notes]: {notes_text}")
        text = "\n".join(parts).strip()
        print(f"  Slide {slide_num + 1}: {len(text)} chars — {text[:80]!r}")
        if text:
            pages.append({
                "source": file_path,
                "source_page": f"Slide {slide_num + 1}",   # ← string
                "text": text
            })
    if not pages:
        print("[WARN] No text extracted from PPT — file may be image-only slides")
    return pages

def load_doc(file_path: str) -> list[dict]:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return open_and_read_pdf(file_path)
    elif ext in (".xlsx", ".xls", ".xlsm"):
        return open_and_read_excel(file_path)
    elif ext in (".pptx", ".ppt"):
        return open_and_read_pptx(file_path)  
    else:
        return open_and_read_document(file_path)


def compute_hash(texts):
    return hashlib.sha256(
        "\n".join(t["text"] for t in texts).encode()
    ).hexdigest()



def _split_into_sentences(text: str) -> list[str]:

    raw = re.split(r'(?<=[.!?])\s+', text.strip())
    return [s.strip() for s in raw if s.strip()]


def semantic_chunk_page(page: dict) -> list[dict]:
    sentences = _split_into_sentences(page["text"])
    if not sentences:
        return []

    embeddings = EMBED_MODEL.encode(sentences, show_progress_bar=False)

    chunks, current_sentences = [], [sentences[0]]

    for i in range(1, len(sentences)):
        dist = cosine_distance(embeddings[i - 1], embeddings[i])
        current_text = " ".join(current_sentences)

        if dist > SEMANTIC_BREAKPOINT_THRESHOLD or len(current_text) > MAX_CHUNK_TOKENS:
            if len(current_text) >= MIN_CHUNK_CHARS:
                chunks.append(current_text)
            current_sentences = [sentences[i]]
        else:
            current_sentences.append(sentences[i])

    tail = " ".join(current_sentences)
    if len(tail) >= MIN_CHUNK_CHARS:
        chunks.append(tail)

    return [
        {
            "text": chunk,
            "source": page["source"],
            "source_page": page["source_page"],
            "chunk_id": None,   # assigned by caller
        }
        for chunk in chunks
    ]


def chunk_pages(texts: list[dict]) -> list[dict]:
    all_chunks, chunk_id = [], 0
    for page in tqdm(texts, desc="Semantic chunking"):
        for c in semantic_chunk_page(page):
            c["chunk_id"] = chunk_id
            all_chunks.append(c)
            chunk_id += 1
    return all_chunks


def generate_embeddings(chunks: list[dict]) -> list[dict]:
    texts = [c["text"] for c in chunks]
    embeddings = EMBED_MODEL.encode(texts, show_progress_bar=True)
    for i, c in enumerate(chunks):
        c["embedding"] = embeddings[i].tolist()
    return chunks


def embed_query(query: str) -> list[float]:
    return EMBED_MODEL.encode([query])[0].tolist()


def get_or_create_collection(name: str):
    try:
        return chroma_client.get_collection(name)
    except Exception:
        return chroma_client.create_collection(name)


def store_chunks(chunks: list[dict], collection, doc_hash: str):
    ids, embs, metas, docs = [], [], [], []
    for c in chunks:
        chunk_id = f"{doc_hash}_{c['chunk_id']}"
        ids.append(chunk_id)
        embs.append(c["embedding"])
        docs.append(c["text"])
        metas.append({
            "source": c["source"],
            "source_page": c["source_page"],
            "chunk_id": c["chunk_id"],
            "doc_hash": doc_hash
        })
    collection.add(ids=ids, embeddings=embs, documents=docs, metadatas=metas)



class BM25Index:
   
    def __init__(self, collection):
        result = collection.get(include=["documents", "metadatas"])
        self.ids       = result["ids"]
        self.docs      = result["documents"]
        self.metas     = result["metadatas"]
        tokenized      = [doc.lower().split() for doc in self.docs]
        self.bm25      = BM25Okapi(tokenized)

    def query(self, query: str, top_k: int) -> list[tuple[str, str, dict, float]]:
    
        tokens = query.lower().split()
        scores = self.bm25.get_scores(tokens)
        top_idx = np.argsort(scores)[::-1][:top_k]
        return [
            (self.ids[i], self.docs[i], self.metas[i], float(scores[i]))
            for i in top_idx
        ]


def _dense_retrieve(query: str, collection, top_k: int):
    q_emb = embed_query(query)
    raw = collection.query(
        query_embeddings=[q_emb],
        n_results=top_k,
        include=["documents", "metadatas", "distances"]
    )
    chunk_ids = raw["ids"][0]
    docs      = raw["documents"][0]
    metas     = raw["metadatas"][0]
    distances = raw["distances"][0]
    return list(zip(chunk_ids, docs, metas, distances))


def _reciprocal_rank_fusion(
    dense_hits:  list[tuple], 
    sparse_hits: list[tuple],
    k: int = RRF_K,
) -> list[tuple]:

    scores: dict[str, float] = {}
    id_to_payload: dict[str, tuple] = {}

    for rank, (cid, doc, meta, _) in enumerate(dense_hits):
        scores[cid]       = scores.get(cid, 0.0) + (1 - BM25_WEIGHT) / (k + rank + 1)
        id_to_payload[cid] = (cid, doc, meta)

    for rank, (cid, doc, meta, _) in enumerate(sparse_hits):
        scores[cid]       = scores.get(cid, 0.0) + BM25_WEIGHT / (k + rank + 1)
        id_to_payload[cid] = (cid, doc, meta)

    ranked = sorted(scores.items(), key=lambda x: x[1], reverse=True)
    return [(id_to_payload[cid][0], id_to_payload[cid][1],
             id_to_payload[cid][2], score)
            for cid, score in ranked]


def _rerank(query: str, candidates: list[tuple], top_k: int) -> list[tuple]:
   
    if not candidates:
        return []

    pairs  = [(query, c[1]) for c in candidates]
    scores = RERANK_MODEL.predict(pairs)           

    ranked = sorted(
        zip(scores, candidates),
        key=lambda x: x[0],
        reverse=True
    )
    return [item for _, item in ranked[:top_k]]

_bm25_cache: dict[str, BM25Index] = {} 

def retrieve(query: str, collection, top_k: int = DEFAULT_TOP_K):
    dense_hits  = _dense_retrieve(query, collection, DENSE_TOP_K)
    if collection.name not in _bm25_cache:         
        _bm25_cache[collection.name] = BM25Index(collection)
    sparse_hits = _bm25_cache[collection.name].query(query, SPARSE_TOP_K)
    fused       = _reciprocal_rank_fusion(dense_hits, sparse_hits)
    reranked    = _rerank(query, fused, top_k)
 
    docs   = [r[1] for r in reranked]
    metas  = [r[2] for r in reranked]
    scores = [r[3] for r in reranked]
 
    return docs, metas, scores



def build_prompt(question: str, docs: list[str]) -> str:
    context = "\n\n".join(docs)
    return f"""You are an expert document analyst. Your job is to give thorough, accurate answers based strictly on the provided context.

INSTRUCTIONS:
- Answer in clear, flowing prose unless the question explicitly asks for a list.
- Be specific — include exact details, numbers, names, conditions, or exceptions mentioned in the context.
- If multiple aspects are relevant, address all of them.
- Explain the "why" or "how" when the context supports it, not just the "what".
- If the context only partially answers the question, answer what you can and clearly state what is not covered.
- If the answer is genuinely not in the context, say: "Not found in the document."
- Never guess, infer, or use outside knowledge.

TABLES:
- If the context contains a table (markdown or plain text), read it row by row and column by column before answering.
- Identify what each column and row represents before drawing conclusions.
- When asked to explain a table, describe: (1) what the table measures or compares, (2) the column/row structure, (3) the most significant values or patterns, and (4) any notable differences, trends, or outliers.
- Always quote specific cell values (e.g. numbers, labels) when they are relevant to the question.
- Do not summarise a table so broadly that specific values are lost.

IMAGES / FIGURES (extracted via OCR):
- If the context contains OCR-extracted text from a figure, diagram, or chart, reconstruct what the visual likely shows based on the extracted labels, values, and captions.
- Describe the type of visual (bar chart, architecture diagram, equation, etc.) if it can be inferred.
- Explain what the visual is communicating, not just what text was found in it.
- If the OCR text is fragmented or unclear, say so and explain what can still be reasonably understood.

CONTEXT:
{context}

QUESTION:
{question}

ANSWER (be thorough and specific):"""


def _generate_with_ollama(prompt: str) -> str:
    try:
        response = requests.post(
            OLLAMA_URL,
            json={
                "model": OLLAMA_MODEL,
                "prompt": prompt,
                "stream": False,
                "options": {
                    "temperature": 0.2,
                    "num_predict": 2048
                }
            },
            timeout=120   # local model can be slow
        )
        response.raise_for_status()
        return response.json()["response"].strip()
    except Exception as e:
        return f"Both Gemini and Ollama failed. Ollama error: {e}"


import requests

OLLAMA_MODEL = "mistral"
OLLAMA_URL   = "http://localhost:11434/api/generate"


def _generate_with_ollama(prompt: str) -> str:
    """Fallback: call local Ollama Mistral model."""
    try:
        response = requests.post(
            OLLAMA_URL,
            json={
                "model": OLLAMA_MODEL,
                "prompt": prompt,
                "stream": False,
                "options": {
                    "temperature": 0.2,
                    "num_predict": 2048
                }
            },
            timeout=120   # local model can be slow
        )
        response.raise_for_status()
        return response.json()["response"].strip()
    except Exception as e:
        return f"Both Gemini and Ollama failed. Ollama error: {e}"


def generate_answer(prompt: str, retries: int = 3, delay: float = 2.0) -> str:
    """Try Gemini with retries, fall back to Ollama on 500."""
    last_error = None

    for attempt in range(retries):
        try:
            response = genai_client.models.generate_content(
                model=LLM_MODEL,
                contents=prompt,
                config=types.GenerateContentConfig(
                    temperature=0.2,
                    max_output_tokens=2048
                )
            )
            return response.text.strip()

        except Exception as e:
            err = str(e)
            last_error = err

            if "500" in err or "INTERNAL" in err:
                if attempt < retries - 1:
                    print(f"[WARN] Gemini 500 — retrying in {delay}s (attempt {attempt + 1}/{retries})")
                    time.sleep(delay)
                    delay *= 2   
                else:
                    print(f"[WARN] Gemini failed after {retries} retries → falling back to Ollama ({OLLAMA_MODEL})")
                    return _generate_with_ollama(prompt)
            else:
                print(f"[WARN] Gemini error ({err}) → falling back to Ollama ({OLLAMA_MODEL})")
                return _generate_with_ollama(prompt)
    print("[WARN] Unexpected fallback to Ollama")
    return _generate_with_ollama(prompt)


def rag_query(question: str, collection, top_k: int = DEFAULT_TOP_K) -> dict:
    if not question.strip():
        return {"answer": "Invalid question", "sources": []}

    docs, metas, scores = retrieve(question, collection, top_k)

    if not docs:
        return {"answer": "No relevant context found", "sources": []}

    prompt = build_prompt(question, docs)
    answer = generate_answer(prompt)

    sources = [
        {
            "source":   os.path.basename(m["source"]),  
            "page":     m["source_page"],                
            "score":    round(s, 4),
            "text":     d
        }
        for d, m, s in zip(docs, metas, scores)
    ]

    return {"answer": answer, "sources": sources}

def ingest_document(file_path: str, original_name: str = None, collection_name: str = COLLECTION_NAME):
    print("Loading document...")
    texts = load_doc(file_path)

    # Replace temp path with original filename if provided
    display_name = original_name or os.path.basename(file_path)
    for t in texts:
        t["source"] = display_name

    doc_hash = compute_hash(texts)          

    print("Semantic chunking...")
    chunks = chunk_pages(texts)
    print(f"  → {len(chunks)} chunks produced")

    print("Embedding...")
    chunks = generate_embeddings(chunks)

    print("Storing in ChromaDB...")
    collection = get_or_create_collection(collection_name)
    for c in chunks:
        c["doc_hash"] = doc_hash        
    store_chunks(chunks, collection, doc_hash)

    print("✓ Ingestion complete\n")
    return collection, len(chunks)


correctness_instructions = """You are a teacher grading a quiz.
You will be given a QUESTION, the GROUND TRUTH (correct) ANSWER, and the STUDENT ANSWER.
Grade the student answers based ONLY on their factual accuracy relative to the ground truth answer.
Ensure that the student answer does not contain any conflicting statements.
It is OK if the student answer contains more information than the ground truth answer, as long as it is factually accurate.
Correctness: True means the student's answer meets all criteria. False means it does not.
Explain your reasoning step-by-step. Avoid simply stating the correct answer at the outset."""

relevance_instructions = """You are a teacher grading a quiz.
You will be given a QUESTION and a STUDENT ANSWER.
Ensure the STUDENT ANSWER is concise and relevant to the QUESTION and helps answer it.
Relevance: True means the student's answer meets all criteria. False means it does not.
Explain your reasoning step-by-step."""

grounded_instructions = """You are a teacher grading a quiz.
You will be given FACTS and a STUDENT ANSWER.
Ensure the STUDENT ANSWER is grounded in the FACTS and does not contain hallucinated information.
Grounded: True means the student's answer meets all criteria. False means it does not.
Explain your reasoning step-by-step."""

retrieval_relevance_instructions = """You are a teacher grading a quiz.
You will be given a QUESTION and a set of FACTS provided by the student.
Identify FACTS that are completely unrelated to the QUESTION.
If the facts contain ANY keywords or semantic meaning related to the question, consider them relevant.
Relevance: True means the FACTS contain ANY keywords or semantic meaning related to the QUESTION. False means they are completely unrelated.
Explain your reasoning step-by-step."""


def gemini_grade(system: str, user: str, bool_key: str) -> bool:
    prompt = f"""{system}

You MUST respond with ONLY valid JSON. No explanation outside JSON.
Use this exact schema:
{{
  "explanation": "<your step-by-step reasoning>",
  "{bool_key}": true or false
}}

{user}"""

    response = genai_client.models.generate_content(
        model=LLM_MODEL,
        contents=prompt,
        config=types.GenerateContentConfig(temperature=0.0, max_output_tokens=512)
    )

    raw = response.text.strip()
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)

    try:
        parsed = json.loads(raw)
        return bool(parsed[bool_key])
    except Exception as e:
        print(f"[WARN] Gemini grader parse error: {e}\nRaw: {raw}")
        return False


def evaluate_rag(collection, test_cases, dataset_name="RAG Evaluation", experiment_prefix="rag-eval"):
    os.environ["LANGSMITH_TRACING"] = "true"

    @traceable()
    def rag_bot(question: str) -> dict:
        result = rag_query(question=question, collection=collection, top_k=3)
        return {"answer": result["answer"], "documents": result["sources"]}

    ls_client = Client()
    if ls_client.has_dataset(dataset_name=dataset_name):
        dataset = ls_client.read_dataset(dataset_name=dataset_name)
    else:
        dataset = ls_client.create_dataset(dataset_name=dataset_name)

    existing_examples = list(ls_client.list_examples(dataset_id=dataset.id))
    if not existing_examples:
        ls_client.create_examples(
            dataset_id=dataset.id,
            examples=[
                {"inputs": {"question": tc["question"]}, "outputs": {"answer": tc.get("ground_truth")}}
                for tc in test_cases
            ],
        )

    def correctness(inputs, outputs, reference_outputs) -> bool:
        return gemini_grade(
            system=correctness_instructions,
            user=(f"QUESTION: {inputs['question']}\n"
                  f"GROUND TRUTH ANSWER: {reference_outputs['answer']}\n"
                  f"STUDENT ANSWER: {outputs['answer']}"),
            bool_key="correct"
        )

    def relevance(inputs, outputs) -> bool:
        return gemini_grade(
            system=relevance_instructions,
            user=(f"QUESTION: {inputs['question']}\n"
                  f"STUDENT ANSWER: {outputs['answer']}"),
            bool_key="relevant"
        )

    def groundedness(inputs, outputs) -> bool:
        facts = "\n\n".join(doc["text"] for doc in outputs["documents"])
        return gemini_grade(
            system=grounded_instructions,
            user=(f"FACTS:\n{facts}\n\nSTUDENT ANSWER: {outputs['answer']}"),
            bool_key="grounded"
        )

    def retrieval_relevance(inputs, outputs) -> bool:
        facts = "\n\n".join(doc["text"] for doc in outputs["documents"])
        return gemini_grade(
            system=retrieval_relevance_instructions,
            user=(f"FACTS:\n{facts}\n\nQUESTION: {inputs['question']}"),
            bool_key="relevant"
        )

    def target(inputs: dict) -> dict:
        return rag_bot(inputs["question"])

    results = ls_client.evaluate(
        target,
        data=dataset_name,
        evaluators=[correctness, groundedness, relevance, retrieval_relevance],
        experiment_prefix=experiment_prefix,
        metadata={"llm": LLM_MODEL, "embed_model": EMBED_MODEL_NAME},
    )


def question(file_path: str):
    original_name = os.path.basename(file_path)   # ← extract filename from path
    collection, _ = ingest_document(file_path, original_name=original_name)
    print("\nReady for Q&A\n")

    while True:
        q = input("Q: ").strip()
        if not q:
            continue

        res = rag_query(q, collection)
        print(f"\nA: {res['answer']}\n")

        print("Sources:")
        for i, s in enumerate(res["sources"], 1):
            print(f"  {i}. {s['source']} — {s['page']} (score: {s['score']})")
        print()

        if input("Continue? (y/n): ").lower() in ["n", "no"]:
            break

if __name__ == "__main__":
    question(r"C:\RAG\Document-Intelligence-RAG\legal.pdf")
'''
    collection = ingest_document("C:\RAG\Document-Intelligence-RAG\legal.pdf", "attention-eval")

    evaluate_rag(
        collection=collection,
        dataset_name="Legal Document RAG Eval",
        experiment_prefix="gemma-attention-rag",
        test_cases = [
    {
        "question": "What happens if a cause of action is not brought within one year?",
        "ground_truth": "It will be deemed forever waived and barred."
    },
    {
        "question": "Can users assign their rights under the Terms of Use?",
        "ground_truth": "No, users may not assign their rights, and any attempt to do so is null and void."
    },
    {
        "question": "Can the company assign its rights under the Terms of Use?",
        "ground_truth": "Yes, the company may freely assign its rights without user consent."
    },
    {
        "question": "What types of information does the site collect?",
        "ground_truth": "The site collects Personally Identifiable Information and Nonpersonally Identifiable Information."
    },
    {
        "question": "What is Personally Identifiable Information?",
        "ground_truth": "It is information that identifies who a person is, such as name, email, address, or other personal details."
    },
    {
        "question": "What is Nonpersonally Identifiable Information?",
        "ground_truth": "It is information that does not identify a specific individual, such as browser type, IP address, and URLs visited."
    },
    {
        "question": "Does the company sell or rent Personally Identifiable Information?",
        "ground_truth": "No, the company does not sell, trade, or rent Personally Identifiable Information."
    },
    {
        "question": "What is an IP address?",
        "ground_truth": "An IP address is a number assigned to a device that allows it to communicate on the Internet."
    },
    {
        "question": "Why does the site collect user information?",
        "ground_truth": "The site collects user information to provide a smooth, efficient, and customized user experience."
    },
    {
        "question": "When is Personally Identifiable Information collected?",
        "ground_truth": "It is collected when users voluntarily provide it during activities like registration, purchases, or surveys."
    },
    {
        "question": "How can users update their personal information?",
        "ground_truth": "Users can update their information by accessing their account online or contacting the company directly."
    },
    {
        "question": "Can users completely delete their information?",
        "ground_truth": "Not always, as some information may remain due to transaction records or backups."
    },
    {
        "question": "What happens if users reject cookies?",
        "ground_truth": "Some parts of the site may not function properly."
    },
    {
        "question": "How is user information protected?",
        "ground_truth": "User information is protected through encryption, secure servers, and restricted access."
    },
    {
        "question": "Does the company guarantee complete data security?",
        "ground_truth": "No, data transmission over the Internet cannot be guaranteed to be 100% secure."
    },
    {
        "question": "Who can access user information within the company?",
        "ground_truth": "Only selected personnel and contractors with password access can access user information."
    },
    {
        "question": "Under what circumstances can user information be shared with authorities?",
        "ground_truth": "User information may be shared when required by law enforcement or judicial authorities through court orders, subpoenas, or warrants."
    },
    {
        "question": "Why might user information be shared with service partners?",
        "ground_truth": "User information may be shared to operate services like payments, delivery, and customer support."
    },
    {
        "question": "What responsibility does a user have regarding their password?",
        "ground_truth": "Users must keep their password secure and not share it with others."
    },
    {
        "question": "What is the minimum age required to use the site?",
        "ground_truth": "Users must be at least 18 years old."
    },
    {
        "question": "What happens if the company collects information from a child under 13?",
        "ground_truth": "The information will be deleted upon discovery."
    },
    {
        "question": "What is the main purpose of cookies on the site?",
        "ground_truth": "Cookies are used to facilitate usage, improve functionality, and personalize user experience."
    },
    {
        "question": "What limitation exists regarding third-party websites?",
        "ground_truth": "The company is not responsible for the privacy policies or practices of third-party websites."
    },
    {
        "question": "What happens if part of the Terms of Use is found invalid?",
        "ground_truth": "The remaining provisions will remain in full force and effect."
    },
    {
        "question": "When do changes to the Privacy Policy become effective?",
        "ground_truth": "Changes become effective immediately upon notice, such as posting or email notification."
    }
]
        
    )
    '''